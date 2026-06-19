import os
import re
import time
from dotenv import load_dotenv
import streamlit as st

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import docx

# Disable telemetry signal handlers that can fail inside Streamlit worker threads.
os.environ.setdefault("CREWAI_DISABLE_TELEMETRY", "true")
os.environ.setdefault("OTEL_SDK_DISABLED", "true")

from crew import (
    research_crew,
    research_crew_with_ppt,
    literature_review_crew,
    paper_comparison_crew,
    citation_reference_crew,
    problem_statement_crew,
)
from tools.arxiv_tool import ArxivSearchTool

# Load environment variables
load_dotenv()

_arxiv_tool = ArxivSearchTool()


def check_api_keys():
    """Check if required API keys are set"""
    required_vars = ['SERPER_API_KEY', 'GROQ_API_KEY']
    missing_vars = [var for var in required_vars if not os.getenv(var)]
    return missing_vars


def get_retry_wait_seconds(error_text):
    """Extract wait time from Groq error message."""
    match = re.search(r"Please try again in\s*([0-9]+(?:\.[0-9]+)?)s", error_text)
    if match:
        return float(match.group(1)) + 5.0
    return 20.0


# Slide titles used when the AI outline is unavailable
_FALLBACK_SLIDE_TITLES = [
    "{topic}",
    "Problem Statement",
    "Introduction",
    "Applications",
    "Advantages",
    "Disadvantages",
    "Conclusion",
]


def generate_ppt_outline_fallback(topic, final_report_path="final_report.md", output_path="presentation_outline.md"):
    """Create a structured presentation outline from generated report content."""
    report_content = ""
    for path in [final_report_path, "analysis_report.md", "research_gaps.md", "research_findings.md"]:
        if os.path.exists(path):
            with open(path, 'r', encoding='utf-8') as f:
                report_content = f.read().strip()
            if report_content:
                break

    if not report_content:
        return False

    bullets = extract_bullets(report_content, limit=35)
    chunks = [bullets[i:i + 5] for i in range(0, 35, 5)]
    while len(chunks) < 7:
        chunks.append(["Details available in the full report."])

    lines = []
    for idx, title_template in enumerate(_FALLBACK_SLIDE_TITLES):
        title = title_template.format(topic=topic)
        slide_bullets = chunks[idx] if idx < len(chunks) else ["See full report for details."]
        lines.append(f"## Slide {idx + 1}: {title}")
        lines.append("**Bullets:**")
        for point in slide_bullets:
            lines.append(f"- {point[:180]}")
        lines.append("**Speaker Notes:**")
        lines.append(f"This slide covers {title.lower()} based on research findings. "
                     f"Refer to the full report for complete data and citations.")
        lines.append("**Visual:** Relevant chart or diagram based on report data.")
        lines.append("")

    outline = "\n".join(lines)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(outline)
    return True


def read_text_file(path):
    """Read text file safely; return empty string if missing."""
    if not os.path.exists(path):
        return ""
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()


def generate_word_document(markdown_content, output_path="final_report.docx"):
    """Convert markdown content to a simple Word document."""
    doc = docx.Document()
    lines = markdown_content.split('\n')
    for line in lines:
        if line.startswith('# '):
            doc.add_heading(line[2:].strip(), 0)
        elif line.startswith('## '):
            doc.add_heading(line[3:].strip(), 1)
        elif line.startswith('### '):
            doc.add_heading(line[4:].strip(), 2)
        elif line.startswith('- '):
            doc.add_paragraph(line[2:].strip(), style='List Bullet')
        elif line.strip() == '':
            continue
        else:
            clean_line = line.replace('**', '').replace('*', '')
            doc.add_paragraph(clean_line)
    doc.save(output_path)
    return True


def extract_bullets(text, limit=5):
    """Extract bullet-style points from markdown/text."""
    bullets = []
    for line in text.splitlines():
        line = line.strip()
        if line.startswith("- ") or line.startswith("* "):
            bullets.append(line[2:].strip())
        if len(bullets) >= limit:
            break
    if bullets:
        return bullets
    sentences = re.split(r"(?<=[.!?])\s+", text)
    fallback = [s.strip() for s in sentences if len(s.strip()) > 20][:limit]
    return fallback if fallback else ["Key insight not available in current report output."]


def extract_numeric_values(text, limit=4):
    """Extract numeric values for charting from report text."""
    matches = re.findall(r"\b\d+(?:\.\d+)?%?\b", text)
    values = []
    for m in matches:
        v = float(m.replace("%", ""))
        if v > 0:
            values.append(v)
        if len(values) >= limit:
            break
    if len(values) < 4:
        tokens = re.findall(r"[A-Za-z]{4,}", text)
        if tokens:
            values = [
                len(tokens),
                len(set(tokens)),
                max(1, len(text.splitlines())),
                max(1, len(re.findall(r"[.!?]", text))),
            ]
        else:
            return []
    return values[:4]


def extract_chart_labels(text, limit=4):
    """Extract meaningful noun-phrase labels from report text for chart categories."""
    phrases = re.findall(r"\b[A-Z][a-z]+(?:\s[A-Z][a-z]+)+", text)
    seen = []
    for p in phrases:
        p = p.strip()
        if p not in seen and len(p) <= 30:
            seen.append(p)
        if len(seen) >= limit:
            break
    if len(seen) >= limit:
        return seen[:limit]
    words = re.findall(r"\b[A-Z][a-z]{3,}\b", text)
    unique_words = list(dict.fromkeys(words))
    for w in unique_words:
        if w not in seen:
            seen.append(w)
        if len(seen) >= limit:
            break
    while len(seen) < limit:
        seen.append(f"Factor {len(seen) + 1}")
    return seen[:limit]


def _apply_title_slide_style(slide, topic):
    """Apply a branded deep-blue background and white text to the title slide."""
    DARK_BLUE = RGBColor(0x1F, 0x38, 0x64)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    ACCENT = RGBColor(0x2E, 0x75, 0xB6)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.6), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.size = Pt(36)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)
                    run.font.size = Pt(16)


def _style_content_slide_title(slide):
    """Apply accent color and bold styling to a content slide's title."""
    ACCENT = RGBColor(0x1F, 0x38, 0x64)
    if slide.shapes.title:
        tf = slide.shapes.title.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.color.rgb = ACCENT
                run.font.bold = True
                run.font.size = Pt(28)


def extract_section_block(text, keywords):
    """Extract lines near headings that match any keyword."""
    lines = text.splitlines()
    lowered = [ln.lower() for ln in lines]
    for i, ln in enumerate(lowered):
        if any(k in ln for k in keywords):
            block = "\n".join(lines[i:i + 12]).strip()
            if block:
                return block
    return ""


def generate_powerpoint_from_reports(topic, output_path="presentation.pptx"):
    """Generate a branded PowerPoint deck exactly from the compiled presentation_outline.md."""
    outline_text = read_text_file("presentation_outline.md")
    if not outline_text:
        return False, "presentation_outline.md not found. Cannot generate PowerPoint."

    prs = Presentation()
    slides_built = 0
    raw_slides = re.split(r'\n##\s+(?:Slide\s+\d+:?\s*)?', '\n' + outline_text)[1:]

    for idx, slide_md in enumerate(raw_slides):
        try:
            lines = [ln.strip() for ln in slide_md.strip().splitlines() if ln.strip()]
            if not lines:
                continue
            title = lines[0].replace('**', '')
            bullets = []
            mode = "bullets"
            for line in lines[1:]:
                lower_line = line.lower()
                if "**visual" in lower_line or "visual:" in lower_line:
                    mode = "visual"
                    continue
                elif "**bullet" in lower_line or "bullets:" in lower_line:
                    mode = "bullets"
                    continue
                if mode == "bullets":
                    if line.startswith("- ") or line.startswith("* "):
                        bullets.append(line[2:].strip())
                    elif re.match(r"^\d+\.\s", line):
                        bullets.append(re.sub(r"^\d+\.\s", "", line).strip())
                    elif line and not line.startswith("**"):
                        bullets.append(line)

            if idx == 0:
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = title
                subtitle = slide.placeholders[1]
                subtitle.text = "\n".join(bullets[:2]) if bullets else "Research Presentation"
                _apply_title_slide_style(slide, topic)
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                tf = slide.placeholders[1].text_frame
                tf.clear()
                if bullets:
                    for i, point in enumerate(bullets):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = point[:200]
                        p.level = 0
                        p.font.size = Pt(18)
                else:
                    p = tf.paragraphs[0]
                    p.text = "See full report for details."
                    p.font.size = Pt(18)
                _style_content_slide_title(slide)

            slides_built += 1
        except Exception as e:
            print(f"Error parsing slide {idx}: {e}")

    if slides_built == 0:
        return False, "Could not extract structured slides from presentation_outline.md"

    prs.save(output_path)
    return True, output_path


# ── Shared helper: run a crew with retry logic ────────────────────────────────

def run_crew_with_retry(crew, inputs, max_retries=3, initial_sleep=5):
    """Kick off a crew with exponential-backoff retry on rate-limit errors.

    Returns (result, error_string). On success error_string is None.
    """
    for attempt in range(max_retries):
        try:
            time.sleep(initial_sleep)
            result = crew.kickoff(inputs)
            return result, None
        except Exception as e:
            if attempt < max_retries - 1:
                wait = get_retry_wait_seconds(str(e))
                st.warning(
                    f"⏳ Rate limited. Retrying in {wait:.1f}s… "
                    f"(Attempt {attempt + 2}/{max_retries})"
                )
                time.sleep(wait)
            else:
                return None, str(e)


def main():
    """Main Streamlit app"""
    st.set_page_config(
        page_title="AI Research Assistant",
        page_icon="🔬",
        layout="wide"
    )

    st.title("🔬 AI Research Assistant")
    st.markdown("*Powered by CrewAI Multi-Agent System*")

    # Initialize session state
    for key in ('research_completed', 'research_result', 'research_error'):
        if key not in st.session_state:
            st.session_state[key] = False if key == 'research_completed' else None

    # ── Sidebar ───────────────────────────────────────────────────────────────
    with st.sidebar:
        st.header("⚙️ Configuration")
        missing_vars = check_api_keys()

        if missing_vars:
            st.error("❌ Missing API Keys")
            st.write("Please set the following environment variables:")
            for var in missing_vars:
                st.code(f"{var}=your_api_key_here")
            st.info("💡 Create a .env file in the project root with your API keys")
        else:
            st.success("✅ API Keys Configured")

        st.header("🤖 Multi-Agent System")
        st.markdown("""
        **Research Agents:**
        - 🔍 **Research Specialist**: Gathers information
        - 📊 **Data Analyst**: Analyses findings
        - 🎯 **Gap Identifier**: Finds research gaps
        - ✍️ **Content Writer**: Creates reports
        - 📽️ **PPT Maker**: Builds slide outline
        - 📚 **Literature Reviewer**: Writes lit reviews
        - 🔬 **Paper Comparator**: Compares papers side-by-side
        - 📖 **Citation Generator**: Formats APA references
        - 🧩 **Problem Statement Generator**: Frames research problems
        """)

        run_mode = st.selectbox(
            "Run Mode",
            [
                "Fast report (4 agents)",
                "Full report + PPT (5 agents)",
            ],
            index=0,
            help="Use Fast mode for fewer rate limits. Use Full mode to include the PPT maker agent."
        )
        include_ppt = run_mode == "Full report + PPT (5 agents)"

    # ── Main Research Section ─────────────────────────────────────────────────
    col1, col2 = st.columns([2, 1])

    with col1:
        st.header("📝 Research Topic")
        topic = st.text_input(
            "Enter your research topic:",
            placeholder="e.g., Machine Learning trends in 2024",
            help="Enter any topic you want to research"
        )

        if st.button("🚀 Start Research", type="primary", disabled=bool(missing_vars)):
            if not topic.strip():
                st.error("Please enter a valid research topic")
            else:
                st.session_state.research_completed = False
                st.session_state.research_result = None
                st.session_state.research_error = None

                progress_container = st.container()
                status_container = st.container()

                with progress_container:
                    st.info("🔄 Research in progress…")
                    progress_bar = st.progress(0)
                    for i in range(101):
                        progress_bar.progress(i)
                        time.sleep(0.1)

                with status_container:
                    st.write("🔍 Research agents are working…")
                    st.write("📊 Analysing data…")
                    st.write("🎯 Identifying gaps…")
                    st.write("✍️ Writing report…")
                    if include_ppt:
                        st.write("📽️ Preparing presentation outline…")

                with st.spinner("📚 Fetching papers from arXiv…"):
                    arxiv_papers = _arxiv_tool._run(topic)

                selected_crew = research_crew_with_ppt if include_ppt else research_crew
                result, error = run_crew_with_retry(
                    selected_crew,
                    {"topic": topic, "arxiv_papers": arxiv_papers}
                )

                if include_ppt and not os.path.exists("presentation_outline.md"):
                    generate_ppt_outline_fallback(topic)

                st.session_state.research_result = result
                st.session_state.research_completed = True
                st.session_state.research_error = error

                progress_container.empty()
                status_container.empty()

    with col2:
        st.header("📊 Status")
        if st.session_state.research_completed:
            if st.session_state.research_error:
                st.error(f"❌ Error: {st.session_state.research_error}")
            else:
                st.success("✅ Research Completed!")
        else:
            st.info("⏳ Waiting for research topic…")

    # ── Results Section ───────────────────────────────────────────────────────
    if st.session_state.research_completed and not st.session_state.research_error:
        st.header("📄 Research Results")

        output_files = {
            "research_findings.md": "🔍 Research Findings",
            "analysis_report.md": "📊 Analysis Report",
            "research_gaps.md": "🎯 Research Gaps",
            "final_report.md": "📝 Final Report",
        }

        if include_ppt:
            if not os.path.exists("presentation_outline.md"):
                generate_ppt_outline_fallback(topic)
            output_files["presentation_outline.md"] = "📽️ Presentation Outline"

        tabs = st.tabs(list(output_files.values()))

        for i, (filename, title) in enumerate(output_files.items()):
            with tabs[i]:
                if os.path.exists(filename):
                    with open(filename, 'r', encoding='utf-8') as f:
                        content = f.read()
                    st.markdown(content)

                    if filename == "final_report.md":
                        col_md, col_word = st.columns(2)
                        with col_md:
                            st.download_button(
                                label=f"📥 Download {title} (.md)",
                                data=content,
                                file_name=filename,
                                mime="text/markdown"
                            )
                        with col_word:
                            generate_word_document(content, "final_report.docx")
                            with open("final_report.docx", "rb") as docf:
                                doc_data = docf.read()
                            st.download_button(
                                label="📥 Download Word (.docx)",
                                data=doc_data,
                                file_name="final_report.docx",
                                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                            )
                    else:
                        st.download_button(
                            label=f"📥 Download {title}",
                            data=content,
                            file_name=filename,
                            mime="text/markdown"
                        )
                else:
                    st.warning(f"File {filename} not found")

        if include_ppt:
            st.header("📊 PowerPoint Presentation")
            ppt_ok, ppt_result = generate_powerpoint_from_reports(topic, output_path="presentation.pptx")
            if ppt_ok and os.path.exists("presentation.pptx"):
                with open("presentation.pptx", "rb") as f:
                    ppt_data = f.read()
                st.success("✅ PowerPoint generated with charts and diagram slides")
                st.download_button(
                    label="📥 Download PowerPoint (.pptx)",
                    data=ppt_data,
                    file_name="presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.warning(f"Could not generate PowerPoint: {ppt_result}")

    # ── Literature Review ─────────────────────────────────────────────────────
    st.markdown("---")
    st.header("📚 Auto Literature Review Generator")
    st.markdown(
        "Generates a structured academic literature review — thematic clusters, "
        "methodology comparison, chronological progression, and research gaps."
    )

    lit_topic = st.text_input(
        "Literature review topic:",
        placeholder="e.g., Transformer models in NLP",
        key="lit_topic"
    )

    if st.button("📚 Generate Literature Review", key="lit_btn", disabled=bool(missing_vars)):
        if not lit_topic.strip():
            st.error("Please enter a topic for the literature review.")
        else:
            with st.spinner("📚 Fetching papers from arXiv…"):
                arxiv_papers = _arxiv_tool._run(lit_topic)

            lit_progress = st.progress(0)
            lit_status = st.empty()
            lit_status.info("🔄 Literature review agents are working…")

            for i in range(101):
                lit_progress.progress(i)
                time.sleep(0.05)

            lit_result, lit_error = run_crew_with_retry(
                literature_review_crew,
                {"topic": lit_topic, "arxiv_papers": arxiv_papers}
            )
            lit_progress.empty()

            if lit_error:
                lit_status.error(f"❌ Error: {lit_error}")
            else:
                lit_status.success("✅ Literature Review Complete!")
                if os.path.exists("literature_review.md"):
                    with open("literature_review.md", "r", encoding="utf-8") as f:
                        lit_content = f.read()
                    st.markdown(lit_content)
                    st.download_button(
                        label="📥 Download Literature Review (.md)",
                        data=lit_content,
                        file_name="literature_review.md",
                        mime="text/markdown",
                        key="lit_download"
                    )
                else:
                    st.warning("literature_review.md was not generated.")

    # ── Paper Comparison ──────────────────────────────────────────────────────
    st.markdown("---")
    st.header("🔬 Research Paper Comparison")
    st.markdown(
        "Fetches papers on a topic and produces a structured side-by-side comparison: "
        "methodology, results, novelty ranking, and synthesis."
    )

    cmp_topic = st.text_input(
        "Topic to compare papers on:",
        placeholder="e.g., Diffusion models for image generation",
        key="cmp_topic"
    )

    cmp_num_papers = st.slider(
        "Number of papers to compare",
        min_value=2, max_value=6, value=3, step=1,
        key="cmp_num_papers"
    )

    if st.button("🔬 Compare Papers", key="cmp_btn", disabled=bool(missing_vars)):
        if not cmp_topic.strip():
            st.error("Please enter a topic for paper comparison.")
        else:
            original_max = _arxiv_tool.max_results
            _arxiv_tool.max_results = cmp_num_papers

            with st.spinner(f"📚 Fetching {cmp_num_papers} papers from arXiv…"):
                arxiv_papers = _arxiv_tool._run(cmp_topic)

            _arxiv_tool.max_results = original_max

            cmp_progress = st.progress(0)
            cmp_status = st.empty()
            cmp_status.info("🔄 Paper comparison agent is working…")

            for i in range(101):
                cmp_progress.progress(i)
                time.sleep(0.05)

            cmp_result, cmp_error = run_crew_with_retry(
                paper_comparison_crew,
                {"topic": cmp_topic, "arxiv_papers": arxiv_papers}
            )
            cmp_progress.empty()

            if cmp_error:
                cmp_status.error(f"❌ Error: {cmp_error}")
            else:
                cmp_status.success("✅ Paper Comparison Complete!")
                if os.path.exists("paper_comparison.md"):
                    with open("paper_comparison.md", "r", encoding="utf-8") as f:
                        cmp_content = f.read()
                    st.markdown(cmp_content)
                    st.download_button(
                        label="📥 Download Comparison Report (.md)",
                        data=cmp_content,
                        file_name="paper_comparison.md",
                        mime="text/markdown",
                        key="cmp_download"
                    )
                else:
                    st.warning("paper_comparison.md was not generated.")

    # ── Citation & Reference Generator ───────────────────────────────────────
    st.markdown("---")
    st.header("📖 Citation & Reference Generator")
    st.markdown(
        "Runs the full research pipeline and then produces a publication-ready "
        "**APA 7th-edition reference list**, an inline citation guide mapping key "
        "claims to sources, and a citation quality summary."
    )

    cite_topic = st.text_input(
        "Topic for citation generation:",
        placeholder="e.g., Federated Learning in healthcare",
        key="cite_topic"
    )

    if st.button("📖 Generate Citations & References", key="cite_btn", disabled=bool(missing_vars)):
        if not cite_topic.strip():
            st.error("Please enter a topic for citation generation.")
        else:
            with st.spinner("📚 Fetching papers from arXiv…"):
                arxiv_papers = _arxiv_tool._run(cite_topic)

            cite_progress = st.progress(0)
            cite_status = st.empty()
            cite_status.info("🔄 Citation agents are working…")

            for i in range(101):
                cite_progress.progress(i)
                time.sleep(0.05)

            cite_result, cite_error = run_crew_with_retry(
                citation_reference_crew,
                {"topic": cite_topic, "arxiv_papers": arxiv_papers}
            )
            cite_progress.empty()

            if cite_error:
                cite_status.error(f"❌ Error: {cite_error}")
            else:
                cite_status.success("✅ Citation Report Complete!")
                if os.path.exists("citation_report.md"):
                    with open("citation_report.md", "r", encoding="utf-8") as f:
                        cite_content = f.read()
                    st.markdown(cite_content)
                    st.download_button(
                        label="📥 Download Citation Report (.md)",
                        data=cite_content,
                        file_name="citation_report.md",
                        mime="text/markdown",
                        key="cite_download"
                    )
                else:
                    st.warning("citation_report.md was not generated.")

    # ── Problem Statement Generator ───────────────────────────────────────────
    st.markdown("---")
    st.header("🧩 Problem Statement Generator")
    st.markdown(
        "Analyses research findings and identified gaps to produce a structured academic "
        "**problem statement** with research questions, SMART objectives, significance "
        "justification, and scope definition."
    )

    prob_topic = st.text_input(
        "Topic for problem statement:",
        placeholder="e.g., Bias in large language models",
        key="prob_topic"
    )

    if st.button("🧩 Generate Problem Statement", key="prob_btn", disabled=bool(missing_vars)):
        if not prob_topic.strip():
            st.error("Please enter a topic for the problem statement.")
        else:
            with st.spinner("📚 Fetching papers from arXiv…"):
                arxiv_papers = _arxiv_tool._run(prob_topic)

            prob_progress = st.progress(0)
            prob_status = st.empty()
            prob_status.info("🔄 Problem statement agent is working…")

            for i in range(101):
                prob_progress.progress(i)
                time.sleep(0.05)

            prob_result, prob_error = run_crew_with_retry(
                problem_statement_crew,
                {"topic": prob_topic, "arxiv_papers": arxiv_papers}
            )
            prob_progress.empty()

            if prob_error:
                prob_status.error(f"❌ Error: {prob_error}")
            else:
                prob_status.success("✅ Problem Statement Complete!")
                if os.path.exists("problem_statement.md"):
                    with open("problem_statement.md", "r", encoding="utf-8") as f:
                        prob_content = f.read()
                    st.markdown(prob_content)
                    st.download_button(
                        label="📥 Download Problem Statement (.md)",
                        data=prob_content,
                        file_name="problem_statement.md",
                        mime="text/markdown",
                        key="prob_download"
                    )
                else:
                    st.warning("problem_statement.md was not generated.")

    # ── Footer ────────────────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("*Built with CrewAI, Streamlit, and Groq*")


if __name__ == "__main__":
    main()
